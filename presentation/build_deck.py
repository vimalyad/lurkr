#!/usr/bin/env python3
"""Build the Lurkr pitch/architecture deck as both PPTX (editable) and HTML
(for a Chrome headless -> PDF render). Content is single-sourced in SLIDES."""
import base64, html, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
LOGO = os.path.join(ROOT, "assets", "icon-only.png")

# Brand palette (matches the app)
BG     = RGBColor(0x0A, 0x0A, 0x0B)
PANEL  = RGBColor(0x16, 0x16, 0x18)
INK    = RGBColor(0xEC, 0xE9, 0xE2)
SIGNAL = RGBColor(0xF5, 0xB5, 0x44)
MUTED  = RGBColor(0x8D, 0x8A, 0x82)
SERIF  = "Georgia"
SANS   = "Helvetica Neue"
MONO   = "Consolas"

SLIDES = [
    {"type": "title", "title": "Lurkr",
     "subtitle": "always watching, never blinking",
     "tagline": "Multi-agent market intelligence — the competitor-watching team that never sleeps."},

    {"kicker": "01 — Problem", "title": "Founders fly blind on competition", "bullets": [
        "Competitor research is manual, scattered across tabs — and stale the day it's done.",
        "Generic “market reports” aren't about your product or your specific edge.",
        "By the time you notice a rival's move, the window to react has already closed."]},

    {"kicker": "02 — Solution", "title": "Describe your idea. Get a personalised brief.", "bullets": [
        "You describe your startup in a sentence; Lurkr finds the real competitors.",
        "A team of AI agents gathers live signals and analyses them for YOUR product.",
        "Output: the single biggest threat + opportunity, each with a recommended action.",
        "And it keeps watching — automatic daily refresh while you sleep."]},

    {"kicker": "03 — How it works", "title": "A grounded multi-agent pipeline", "bullets": [
        "①  Discovery — finds real, currently-operating competitors in your space.",
        "②  Gather — live web (Tavily) + news (Google News) signals per competitor.",
        "③  Analysts — Marketing · Product · Sales run in parallel, grounded in those signals.",
        "④  Strategy — synthesises one executive brief: threat + opportunity + watch items."]},

    {"kicker": "04 — Architecture", "title": "Tiered, secret-isolating topology", "bullets": [
        "Thin-client Android shell (Capacitor) loads a CDN-hosted React 19 SPA.",
        "Stateless Express API tier holds every key — nothing secret ships on the client.",
        "Neon Postgres stores accounts, ideas and cached analyses.",
        "OpenRouter (LLMs) · Tavily + Google News (signals) · GitHub Actions (cron + CI/CD)."]},

    {"kicker": "05 — The engine", "title": "Why the multi-agent design wins", "bullets": [
        "Fan-out / fan-in: three analysts run concurrently, Strategy synthesises the result.",
        "Grounding / RAG discipline: every claim must cite a real signal — anti-hallucination.",
        "Cost-aware model tiering: a fast model for breadth, a stronger one for judgement.",
        "Fault-isolated sources + bounded-retry JSON parsing — the sweep always completes."]},

    {"kicker": "06 — Stack & rationale", "title": "Why this stack", "bullets": [
        "Neon (serverless Postgres) — relational + JSONB, scale-to-zero, DB branching, zero lock-in.",
        "OpenRouter — one OpenAI-compatible API across models; cost-aware tiering, no lock-in.",
        "Roll-our-own auth — scrypt + HMAC JWT on the Node std-lib; no auth SDK, data stays ours.",
        "Render + GitHub Pages + Actions — free-tier, zero-ops, ship by git push."]},

    {"kicker": "07 — Accounts & memory", "title": "Per-user intelligence that compounds", "bullets": [
        "Hard auth gate — Google (native account picker on Android) or email + password.",
        "Every search is saved per user; a “My Ideas” library keeps them.",
        "Stale-while-present cache: instant re-opens; a live sweep runs only on a cache miss.",
        "Opt-in daily refresh at 04:00 IST — a cron re-runs and caches your watched ideas."]},

    {"kicker": "08 — Delivery", "title": "Build-once, ship by git push", "bullets": [
        "The Android APK is a thin shell — frontend updates need no rebuild or reinstall.",
        "A signed, auto-versioned APK pipeline publishes to GitHub Releases.",
        "Frontend auto-deploys to GitHub Pages; the backend auto-deploys on Render.",
        "A CI build-gate compiles the app on every change before it ships."]},

    {"kicker": "09 — Engineering", "title": "Security & resilience by default", "bullets": [
        "Secret isolation — keys live only on the backend; the client bundle ships none.",
        "Deny-by-default gate; the user is revalidated from the DB on every request.",
        "Memory-hard scrypt hashing, constant-time verify, audience-pinned Google tokens.",
        "Idempotent migrations; an append-only usage ledger ready for usage-based billing."]},

    {"kicker": "10 — Roadmap", "title": "What's next", "bullets": [
        "Usage-based billing — the usage ledger is already in place.",
        "Email alerts when a new high-urgency threat appears (Resend).",
        "Email verification + password reset, once a sending domain is provisioned.",
        "Native Google sign-in polish and an iOS build."]},

    {"type": "closing", "title": "Lurkr",
     "subtitle": "Always watching, never blinking.",
     "tagline": "React 19 · Vite · Tailwind v4 · Capacitor 7 · Express · Neon Postgres · OpenRouter · GitHub Actions · Render"},
]

# ─────────────────────────────── PPTX ───────────────────────────────
def build_pptx(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    def bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

    def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def style(p, text, size, color, bold=False, font=SANS, space_after=10, align=PP_ALIGN.LEFT):
        p.text = text
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.runs[0]
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = font; r.font.color.rgb = color

    def accent(slide, l, t, w=Inches(0.55), h=Inches(0.07)):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = SIGNAL
        shp.line.fill.background()
        shp.shadow.inherit = False

    def footer(slide, n):
        tf = textbox(slide, Inches(0.7), Inches(7.0), Inches(6), Inches(0.4))
        style(tf.paragraphs[0], "LURKR", 9, MUTED, font=MONO, space_after=0)
        tf2 = textbox(slide, Inches(11.6), Inches(7.0), Inches(1.0), Inches(0.4))
        style(tf2.paragraphs[0], f"{n:02d}", 9, MUTED, font=MONO, space_after=0, align=PP_ALIGN.RIGHT)

    for i, s in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank); bg(slide)
        typ = s.get("type", "content")
        if typ in ("title", "closing"):
            if os.path.exists(LOGO):
                slide.shapes.add_picture(LOGO, Inches(5.79), Inches(1.5), height=Inches(1.75))
            tf = textbox(slide, Inches(1), Inches(3.4), Inches(11.33), Inches(2.5), MSO_ANCHOR.TOP)
            style(tf.paragraphs[0], s["title"], 66, INK, bold=True, font=SERIF, space_after=6, align=PP_ALIGN.CENTER)
            p = tf.add_paragraph(); style(p, s["subtitle"], 22, SIGNAL, font=SERIF, space_after=16, align=PP_ALIGN.CENTER)
            p = tf.add_paragraph(); style(p, s["tagline"], 13, MUTED, font=SANS, space_after=0, align=PP_ALIGN.CENTER)
        else:
            accent(slide, Inches(0.75), Inches(0.85))
            tf = textbox(slide, Inches(0.75), Inches(1.0), Inches(11.8), Inches(1.5))
            style(tf.paragraphs[0], s["kicker"], 12, SIGNAL, font=MONO, space_after=4)
            p = tf.add_paragraph(); style(p, s["title"], 36, INK, bold=True, font=SERIF, space_after=0)
            body = textbox(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
            first = True
            for b in s["bullets"]:
                p = body.paragraphs[0] if first else body.add_paragraph()
                first = False
                p.space_after = Pt(16)
                run = p.add_run(); run.text = "▸  "
                run.font.size = Pt(16); run.font.color.rgb = SIGNAL; run.font.name = MONO; run.font.bold = True
                run2 = p.add_run(); run2.text = b
                run2.font.size = Pt(16); run2.font.color.rgb = INK; run2.font.name = SANS
            footer(slide, i + 1)
    prs.save(path)
    return path

# ─────────────────────────────── HTML (for PDF) ───────────────────────────────
def build_html(path):
    logo_b64 = ""
    if os.path.exists(LOGO):
        logo_b64 = base64.b64encode(open(LOGO, "rb").read()).decode()
    css = """
    @page { size: 13.333in 7.5in; margin: 0; }
    * { margin:0; padding:0; box-sizing:border-box; -webkit-print-color-adjust:exact; print-color-adjust:exact; }
    body { font-family:'Helvetica Neue',Arial,sans-serif; color:#ECE9E2; }
    .slide { width:13.333in; height:7.5in; background:#0a0a0b; position:relative;
             padding:0.85in 0.9in; page-break-after:always; overflow:hidden; }
    .slide::before { content:""; position:absolute; inset:0;
      background: radial-gradient(120% 80% at 50% -10%, rgba(245,181,68,.10), transparent 55%),
        linear-gradient(to right, rgba(255,255,255,.022) 1px, transparent 1px),
        linear-gradient(to bottom, rgba(255,255,255,.022) 1px, transparent 1px);
      background-size:100% 100%, 46px 46px, 46px 46px; pointer-events:none; }
    .kicker { font-family:Consolas,monospace; font-size:13px; letter-spacing:.18em; text-transform:uppercase; color:#f5b544; }
    .bar { width:54px; height:6px; background:#f5b544; border-radius:3px; margin-bottom:18px; }
    h1.title { font-family:Georgia,serif; font-size:40px; line-height:1.05; margin-top:8px; max-width:11in; }
    ul { list-style:none; margin-top:34px; }
    li { font-size:19px; line-height:1.5; color:#d9d6cf; margin-bottom:20px; display:flex; gap:14px; max-width:10.6in; }
    li .b { color:#f5b544; font-family:Consolas,monospace; font-weight:700; }
    .footer { position:absolute; bottom:0.55in; left:0.9in; right:0.9in; display:flex; justify-content:space-between;
              font-family:Consolas,monospace; font-size:11px; letter-spacing:.16em; color:#8d8a82; }
    /* hero (title / closing) */
    .hero { display:flex; flex-direction:column; align-items:center; justify-content:center; text-align:center; height:100%; }
    .hero img { width:128px; height:128px; border-radius:28px; box-shadow:0 12px 40px rgba(0,0,0,.5); margin-bottom:26px; }
    .hero h1 { font-family:Georgia,serif; font-size:88px; line-height:.9; }
    .hero .sub { font-family:Georgia,serif; font-style:italic; font-size:26px; color:#f5b544; margin-top:10px; }
    .hero .tag { font-size:15px; color:#8d8a82; margin-top:22px; max-width:8.5in; line-height:1.6; }
    """
    parts = [f"<!doctype html><html><head><meta charset='utf-8'><style>{css}</style></head><body>"]
    for i, s in enumerate(SLIDES):
        typ = s.get("type", "content")
        if typ in ("title", "closing"):
            img = f"<img src='data:image/png;base64,{logo_b64}'/>" if logo_b64 else ""
            parts.append(f"<section class='slide'><div class='hero'>{img}"
                         f"<h1>{html.escape(s['title'])}</h1>"
                         f"<div class='sub'>{html.escape(s['subtitle'])}</div>"
                         f"<div class='tag'>{html.escape(s['tagline'])}</div></div></section>")
        else:
            lis = "".join(f"<li><span class='b'>&#9656;</span><span>{html.escape(b)}</span></li>" for b in s["bullets"])
            parts.append(f"<section class='slide'><div class='kicker'>{html.escape(s['kicker'])}</div>"
                         f"<div class='bar'></div><h1 class='title'>{html.escape(s['title'])}</h1>"
                         f"<ul>{lis}</ul>"
                         f"<div class='footer'><span>LURKR</span><span>{i+1:02d}</span></div></section>")
    parts.append("</body></html>")
    open(path, "w").write("".join(parts))
    return path

if __name__ == "__main__":
    p1 = build_pptx(os.path.join(HERE, "Lurkr-Deck.pptx"))
    p2 = build_html(os.path.join(HERE, "Lurkr-Deck.html"))
    print("wrote", p1)
    print("wrote", p2)
